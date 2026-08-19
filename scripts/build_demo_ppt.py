# -*- coding: utf-8 -*-
"""生成参赛汇报 PPT（作品：智能物资数据管理系统）。

用法: /usr/bin/python3 scripts/build_demo_ppt.py
输出: /workspace/2026-07/智能物资数据管理系统-参赛汇报.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path("/workspace/2026-07/智能物资数据管理系统-参赛汇报.pptx")

BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
LIGHT = RGBColor(0xF2, 0xF6, 0xFA)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)

SW, SH = Inches(13.333), Inches(7.5)


def _bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=LIGHT, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    return sp


def _txt(slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def _bullets(slide, x, y, w, h, items, size=16, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("▪ " if isinstance(it, str) else it[0])
        r.font.size = Pt(size)
        r.font.color.rgb = DARK
        if isinstance(it, tuple):
            r2 = p.add_run()
            r2.text = it[1]
            r2.font.size = Pt(size)
            r2.font.color.rgb = DARK
            r.font.bold = True
    return tb


def _title(slide, text, sub=None):
    _txt(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7), text, size=30, color=BLUE, bold=True)
    if sub:
        _txt(slide, Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.4), sub, size=14, color=GRAY)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.42), Inches(1.6), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()


def _footer(slide, n):
    _txt(slide, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.35), str(n), size=12, color=GRAY, align=PP_ALIGN.RIGHT)


def main() -> int:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    # ---------- 1 封面 ----------
    s = prs.slides.add_slide(blank)
    _bg(s, BLUE)
    _box(s, Inches(0), Inches(5.9), SW, Inches(1.6), fill=RGBColor(0x18, 0x3E, 0x60))
    _txt(s, Inches(1.2), Inches(2.1), Inches(11), Inches(1.0), "智能物资数据管理系统", size=44, color=WHITE, bold=True)
    _txt(s, Inches(1.2), Inches(3.2), Inches(11), Inches(0.6), "基于本地大模型的物资数据治理与智能问答平台", size=20, color=RGBColor(0xBF, 0xD7, 0xEF))
    _txt(s, Inches(1.2), Inches(4.0), Inches(11), Inches(0.5), "AI 辅助开发 · 本地模型推理 · 全程保密合规", size=15, color=RGBColor(0x9F, 0xC4, 0xE8))
    _txt(s, Inches(1.2), Inches(6.15), Inches(11), Inches(1.2), "参赛人：__________   所在单位/部门：__________", size=16, color=WHITE)

    # ---------- 2 背景与痛点 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "背景与痛点", "多源异构台账 · 口径不一 · 查数难")
    _bullets(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.0), [
        ("物资台账分散存放：", "库存、备品备件、需求、资产、出入库流水以 Excel/CSV 多格式零散管理"),
        ("多源异构、口径不一：", "字段定义、日期格式、计量单位、编码体系各不相同"),
        ("人工治理成本高：", "合并清洗耗时费力、易出错、过程不可追溯、结果不可回滚"),
        ("查数依赖人工翻表：", "管理决策所需指标难以快速、口径统一地获取"),
        ("保密要求高：", "涉及内部业务数据，严禁上传公网大模型，必须本地化处理"),
    ], size=18, gap=14)

    # ---------- 3 总体方案 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "总体方案：上传—治理—发布—问答 全链路", "系统架构与数据流转")
    steps = ["① 数据接入\n多 Sheet 解析\n证据留存", "② AI 治理\n画像/质量/映射\n流水拆解·勾稽", "③ 可信发布\n防重写入\n来源可溯·审计·备份", "④ 智能问答\nText2SQL\n指标模板·图表"]
    xs = [0.8, 4.05, 7.3, 10.55]
    for i, st in enumerate(steps):
        _box(s, Inches(xs[i]), Inches(2.0), Inches(2.7), Inches(1.7), fill=BLUE if i != 1 else ACCENT)
        _txt(s, Inches(xs[i] + 0.15), Inches(2.25), Inches(2.4), Inches(1.3), st, size=15, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            _txt(s, Inches(xs[i] + 2.7), Inches(2.45), Inches(0.5), Inches(0.8), "→", size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    _bullets(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.6), [
        ("技术栈：", "FastAPI + Vue3/Element Plus + SQLite（元数据/任务）+ DuckDB（分析库）+ Parquet（证据层）+ vLLM"),
        ("部署：", "Docker Compose 分服务编排；项目组成 app / frontend / deploy / scripts / tests"),
        ("关键设计：", "唯一写入入口、规则+模型双通道、人工确认门槛、全流程审计与一键回滚"),
    ], size=16, gap=10)

    # ---------- 4 核心能力：接入与解析 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "核心能力 ① 数据接入与自动解析", "上传即解析 · 原始证据留存")
    cards = [
        ("多 Sheet 自动识别", "自动解析多个业务 Sheet（维护材料/备品备件/应急备汛/工器具）及其列结构"),
        ("异构格式兼容", "兼容 Excel/CSV、多级表头、换行单元格、混合日期与单位格式"),
        ("原始证据留存", "每次上传保留原始文件快照与 Parquet 证据层，供回溯与审计"),
        ("统一字段画像", "自动生成字段语义画像（分类、占比、样例、单位），支撑下游映射决策"),
    ]
    for i, (t, d) in enumerate(cards):
        x = 0.8 + (i % 2) * 6.15
        y = 2.0 + (i // 2) * 2.1
        _box(s, Inches(x), Inches(y), Inches(5.7), Inches(1.85), fill=LIGHT, line=ACCENT)
        _txt(s, Inches(x + 0.25), Inches(y + 0.2), Inches(5.2), Inches(0.5), t, size=17, color=BLUE, bold=True)
        _txt(s, Inches(x + 0.25), Inches(y + 0.75), Inches(5.2), Inches(1.0), d, size=14, color=DARK)

    # ---------- 5 核心能力：AI 治理 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "核心能力 ② AI 治理", "本地大模型 + 规则引擎双通道")
    _bullets(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.0), [
        ("字段映射建议：", "Embedding 向量召回 + 规则兜底，给出列映射建议，人工一键确认"),
        ("出入库流水智能拆解：", "把“2026.1.6/张伟/会议室搭建备用”类文本自动拆为 时间/人员/数量/用途 结构化流水"),
        ("勾稽差异自动暴露：", "校验“初始+入库−出库 = 现有库存”，口径不一致自动清单化（示例：39 条）"),
        ("数据质量预检：", "空值、格式、单位、编码规范性全量检查，治理前先体检"),
        ("人工确认门槛：", "关键决策必须人工确认，保障准确性可审计、可回滚"),
    ], size=18, gap=14)

    # ---------- 6 核心能力：可信发布 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "核心能力 ③ 可信发布", "唯一入口 · 防重写入 · 来源可溯 · 审计")
    cards = [
        ("唯一写入入口", "所有发布经单一服务路径，杜绝绕过审计的写入"),
        ("防重复发布", "重复确认不会重复入账，发布可安全重试"),
        ("来源可追溯", "标准表行级关联来源文件/Sheet/映射规则，可核查"),
        ("审计时间线", "上传/治理/发布全操作留痕，支持按需回滚与备份"),
    ]
    for i, (t, d) in enumerate(cards):
        x = 0.8 + (i % 2) * 6.15
        y = 2.0 + (i // 2) * 2.1
        _box(s, Inches(x), Inches(y), Inches(5.7), Inches(1.85), fill=LIGHT, line=ACCENT)
        _txt(s, Inches(x + 0.25), Inches(y + 0.2), Inches(5.2), Inches(0.5), t, size=17, color=BLUE, bold=True)
        _txt(s, Inches(x + 0.25), Inches(y + 0.75), Inches(5.2), Inches(1.0), d, size=14, color=DARK)

    # ---------- 7 核心能力：智能问答 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "核心能力 ④ 智能问答", "自然语言问数 · 秒级出数")
    _box(s, Inches(0.8), Inches(1.9), Inches(6.0), Inches(4.6), fill=LIGHT)
    _txt(s, Inches(1.05), Inches(2.1), Inches(5.5), Inches(0.5), "问数示例", size=17, color=BLUE, bold=True)
    _bullets(s, Inches(1.05), Inches(2.7), Inches(5.5), Inches(3.6), [
        ("“各存放位置库存数量排名”", "→ 自动生成 GROUP BY 聚合 SQL"),
        ("“入库和出库分别多少条、合计多少数量”", "→ 流水类型汇总"),
        ("“库存低于最低阈值的物资有哪些”", "→ 补货预警"),
        ("“某水电站工具类库存总量”", "→ 多条件过滤 + 分类聚合"),
    ], size=15, gap=10)
    _box(s, Inches(7.1), Inches(1.9), Inches(5.4), Inches(4.6), fill=LIGHT)
    _txt(s, Inches(7.35), Inches(2.1), Inches(4.9), Inches(0.5), "实现机制", size=17, color=BLUE, bold=True)
    _bullets(s, Inches(7.35), Inches(2.7), Inches(4.9), Inches(3.6), [
        ("Text2SQL：", "本地大模型理解自然语言并生成 SQL"),
        ("指标模板优先：", "高频指标走模板，稳、快、省"),
        ("只读防注入：", "查询限只读并校验，杜绝数据篡改"),
        ("结果可视化：", "表格 + 图表一键切换"),
    ], size=15, gap=10)

    # ---------- 8 模型 API 与保密合规 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "模型 API 接入与保密合规", "本地部署 · 数据不出本机")
    _box(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.4), fill=LIGHT, line=ACCENT)
    _txt(s, Inches(1.05), Inches(2.1), Inches(11.2), Inches(2.0), (
        "生成大模型：Qwen3.6-27B（vLLM，OpenAI 兼容接口 :8001）—— 负责映射建议、流水拆解、Text2SQL\n"
        "向量模型：Qwen3-Embedding-0.6B（vLLM，OpenAI 兼容接口 :8002）—— 负责语义召回/列映射\n"
        "全部模型本地部署、内网回环调用；开发用 AI 平台完成，运行时数据全程不出本机"
    ), size=17, color=DARK)
    _bullets(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.2), [
        ("保密红线：", "严禁内部业务数据上传公网大模型，本方案从架构上杜绝该风险"),
        ("合规治理：", "多模型按场景路由、评测集把关、人工确认门槛、全链路审计"),
        ("演示合规：", "演示使用脱敏样例数据（虚构站点/人员/物料），可一键重建干净演示环境"),
    ], size=17, gap=10)

    # ---------- 9 演示成果 ----------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "演示成果（脱敏样例数据实测）", "一次运行快照")
    nums = [("49", "标准库存行"), ("48", "流水拆解段"), ("39", "勾稽差异发现"), ("22", "流水物料对齐")]
    for i, (n, lab) in enumerate(nums):
        x = 0.8 + i * 3.05
        _box(s, Inches(x), Inches(1.95), Inches(2.75), Inches(1.6), fill=BLUE)
        _txt(s, Inches(x), Inches(2.2), Inches(2.75), Inches(0.7), n, size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, Inches(x), Inches(2.95), Inches(2.75), Inches(0.5), lab, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    _bullets(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(3.0), [
        ("输入：", "4-Sheet 脱敏台账（维护材料/备品备件/应急备汛物资/公用工器具，含脏数据特征）"),
        ("输出：", "标准库存表 fact_inventory + 出入库流水 fact_stock_flow + 勾稽差异清单 + 问数结果"),
        ("问数样例：", "“各存放位置库存数量排名”→ 望湖大厦6楼通信材料室 47；锦澜电站左岸厂房仓库 47 …"),
        ("复现：", "python3 scripts/build_demo_env.py 一键重建干净演示库"),
    ], size=17, gap=12)

    # ---------- 10 价值与展望 ----------
    s = prs.slides.add_slide(blank)
    _bg(s, BLUE)
    _txt(s, Inches(1.2), Inches(2.0), Inches(11), Inches(0.8), "价值与展望", size=32, color=WHITE, bold=True)
    _bullets(s, Inches(1.2), Inches(3.0), Inches(10.9), Inches(3.0), [
        ("治理效率：", "多源异构台账“上传即治理”，人工数周 → 分钟级"),
        ("数据可信：", "标准口径统一、全程来源可溯、可回滚可备份"),
        ("决策赋能：", "自然语言问数，让业务人员零门槛取数"),
        ("复制推广：", "沉淀可复用治理管线，可扩展至需求/资产等其他数据域"),
    ], size=18, gap=14)
    _txt(s, Inches(1.2), Inches(6.4), Inches(11), Inches(0.6), "谢谢观看", size=24, color=RGBColor(0xBF, 0xD7, 0xEF), bold=True)

    prs.save(str(OUT))
    print("已生成:", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
